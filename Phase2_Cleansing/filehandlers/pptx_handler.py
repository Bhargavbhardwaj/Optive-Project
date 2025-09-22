import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE


from ..detectors import detect_pii_in_text
from ..maskers import mask_text
#from ..audit import write_audit_row
from Phase2_Cleansing.audit import AuditLogger
audit = AuditLogger("audit_log.csv")

def clean_pptx_file(input_path, output_path, action, use_spacy,audit, remove_immages = True):
    try:
        prs = pptx.Presentation(input_path) # loads the file into prs object
    except FileNotFoundError:
        print(f"[FAIL] PowerPoint file not found: {input_path}")
        return False
    except Exception as e:
        print(f"[FAIL] Could not open PPTX file {input_path}: {e}")
        return False

    try:
        for slide_idx, slide in enumerate(prs.slides): # goes through every silde
            for shape in slide.shapes: # goes through every shape in the slide
                if hasattr(shape, "text") and shape.text: # if the shape has text and it's not empty then scan it for PII
                    detections = detect_pii_in_text(shape.text, use_spacy=use_spacy)
                    if detections: # if found PII, then mask_text to remove it and store it in cleaned_result
                        cleaned_text = mask_text(shape.text, detections, action=action)

                        # update the shape text and ensure that text really gets updates regardless of shape type
                        try: # tries to set teh shape's text directly and if fails than goes usual method
                            shape.text = cleaned_text
                        except Exception:
                            if hasattr(shape, "text_frame") and shape.has_text_frame: # if fails, then clears all existing para inside the shape's text_frame
                                for para in shape.text_frame.paragraphs:
                                    para.text = ""
                                shape.text_frame.paragraphs[0].text = cleaned_text  # write the cleaned text into the first paragraph
                        for d in detections:  # logs each detection
                            audit.write_row(input_path, output_path, d.get("source"),
                                            d.get("type"), d.get("match"), action, notes=f"slide:{slide_idx+1}")
                  # logs which slide number the detection was found

                # optionally removes image (useful in case image contains sensitive data)
                if remove_immages and shape.shape_type ==13:   # if shape_type ==13, then remove the image from the slide
                    # shape_type 13 means, it is picture shape as per python-pptx
# NOTE :  A PowerPoint file (.pptx) is really a ZIP archive of XML files (following the Office Open XML format).
                    sp = shape._element  # gives you the raw XML element that defines the shape inside the PowerPoint file.
                    sp.getparent().remove(sp)  # tells the parent to remove this child node (sp) from the XML tree
                    audit.write_row(input_path, output_path, "image_removal",
                                    "IMAGE", "image_removed", "remove",
                                    notes = f"slide{slide_idx+1}")

        prs.save(output_path)
        return True
    except Exception as e:
        print(f"[FAIL] Error processing PPTX file {input_path}: {e}")
        return False

