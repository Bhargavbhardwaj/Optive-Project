'''This file is a unified text extraction module that takes different
file types and converts their contents into plain text strings using
the appropriate libraries.'''



import os  # for file path handling
import pytesseract  # pytesseract + PIL.Image - OCR text from images.
from pdfminer.high_level import extract_text as pdf_extract_text  # extract text from pdfs
from pptx import Presentation # read powerpoint slides
import docx # reads docs file
import openpyxl # for excel spreadsheets
from PIL import Image



# opens image with pillow.
def extract_from_image(file_path):
    try:
        text = pytesseract.image_to_string(Image.open(file_path)) # Uses tesseract OCR to detect and extract text.
        return text.strip() # .strip() removes leading whitespace
    except Exception as e:
        return f"[ERROR extracting image: {e}]"

def extract_from_pptx(file_path):
    try:
        prs = Presentation(file_path) # loads a ppt file with Presentation() method
        text = []
        for slides in prs.slides:  # iterates thorugh slides and its shapes
            for shape in slides.shapes:
                if hasattr(shape, "text"):  # extract if it contains text
                    text.append(shape.text)
        return "\n".join(text)  # collects text into a list then joins into a single string.
    except Exception as e:
        return f"[ERROR extracting pptx: {e}]"

def extract_from_pdf(file_path):
    try:
        return pdf_extract_text(file_path) # uses pdfminer.six's extract_text() to read all text from PDF
    except Exception as e:
        return f"[ERROR extracting pdf: {e}]"

def extract_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs]) # extracts all paragraphs’ text and join them with newlines
    except Exception as e:
        return f"[ERROR extracting docx: {e}]"

def extract_from_xlsx(file_path):
    try:
        wb = openpyxl.load_workbook(file_path,data_only=True) # loads the excel file
        text = []
        for sheet in wb.sheetnames:  # reads all sheets - rows - cells
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                # converts each row into a space-separated string
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text.append(row_text)  # collects non-empty rows and joins into final text
        return "\n".join(text)
    except Exception as e:
        return f"[ERROR extracting xlsx: {e}]"

def extract_from_text(file_path):
    try:
        with open(file_path, "r",encoding="utf-8", errors="ignore") as f:
            return f.read()   # reads full file content as string
    except Exception as e:
        return f"[ERROR extracting text: {e}]"


# detects the extension and calls the correct extractors finction based on file type
def extract_content(file_path, file_type):
    ext = file_type.lower().strip(".")
    if ext in ["png", "jpg", "jpeg"]:
        return extract_from_image(file_path)
    elif ext == "pdf":
        return extract_from_pdf(file_path)
    elif ext == "pptx":
        return extract_from_pptx(file_path)
    elif ext == "docx":
        return extract_from_docx(file_path)
    elif ext in ["xlsx","xls"]:
        return extract_from_xlsx(file_path)
    elif ext in ["text", "csv", "log", "json"]:
        return extract_from_text(file_path)
    else:
        return "[Unsupported file type]"


